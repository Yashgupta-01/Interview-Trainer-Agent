from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Color Palette
BG_DARK     = RGBColor(0x0D, 0x11, 0x1C)   # Deep navy
BG_CARD     = RGBColor(0x16, 0x1E, 0x35)   # Card navy
ACCENT_PURP = RGBColor(0x7C, 0x3A, 0xED)   # Purple
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)   # Cyan blue
ACCENT_TEAL = RGBColor(0x10, 0xB9, 0x81)   # Teal green
TEXT_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED  = RGBColor(0x94, 0xA3, 0xB8)
TEXT_YELLOW = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_RED  = RGBColor(0xEF, 0x44, 0x44)

WORKFLOW_IMG = r"C:\Users\hp\.gemini\antigravity\brain\78ee79e2-2f91-4d39-b358-7c4b43760b26\workflow_flowchart_1783753995294.png"
ARCH_IMG     = r"C:\Users\hp\.gemini\antigravity\brain\78ee79e2-2f91-4d39-b358-7c4b43760b26\architecture_blueprint_1783752674967.png"
OUT_PATH     = r"C:\Users\hp\.gemini\antigravity\brain\78ee79e2-2f91-4d39-b358-7c4b43760b26\AI_Interview_Trainer.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # fully blank layout

# ── helpers ──────────────────────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill=BG_CARD, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=TEXT_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_pill(slide, text, l, t, fill=ACCENT_PURP, text_color=TEXT_WHITE, size=11):
    w = len(text) * 0.085 + 0.25
    box = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.3))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color

def slide_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.1, fill=BG_CARD, line=ACCENT_PURP)
    add_text(slide, title, 0.35, 0.12, 10, 0.6, size=28, bold=True, color=ACCENT_BLUE)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.65, 10, 0.4, size=13, color=TEXT_MUTED)

def bullet_card(slide, l, t, w, h, title, bullets, title_color=ACCENT_TEAL, bullet_size=13):
    add_rect(slide, l, t, w, h, fill=BG_CARD, line=ACCENT_PURP)
    add_text(slide, title, l+0.15, t+0.1, w-0.3, 0.35, size=14, bold=True, color=title_color)
    # bullets
    by = t + 0.5
    for b in bullets:
        add_text(slide, f"  •  {b}", l+0.1, by, w-0.2, 0.38, size=bullet_size, color=TEXT_WHITE)
        by += 0.37


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)

# big gradient rect decoration
add_rect(sl, 0, 5.8, 13.33, 1.7, fill=RGBColor(0x1E, 0x12, 0x40))

# glow circle (approximate)
circ = sl.shapes.add_shape(9, Inches(5.5), Inches(1.0), Inches(2.3), Inches(2.3))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x4C, 0x1D, 0x95)
circ.line.fill.background()

add_text(sl, "🤖", 5.95, 1.2, 1.4, 1.0, size=48, align=PP_ALIGN.CENTER)

add_text(sl, "AI Interview Trainer", 1.5, 3.0, 10.3, 1.0,
         size=44, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text(sl, "Powered by IBM Cloud Watsonx AI  ·  FastAPI  ·  ChromaDB  ·  RAG",
         1.5, 4.05, 10.3, 0.5, size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
add_text(sl, "IBM Edunet Summer Internship Project", 1.5, 4.6, 10.3, 0.4,
         size=13, color=TEXT_YELLOW, align=PP_ALIGN.CENTER)

add_pill(sl, "FastAPI", 2.3, 6.1, ACCENT_PURP)
add_pill(sl, "ChromaDB", 3.5, 6.1, ACCENT_BLUE, RGBColor(0,0,0))
add_pill(sl, "IBM Watsonx", 5.0, 6.1, RGBColor(0x17, 0x5C, 0xD6))
add_pill(sl, "LLaMA 3.3 70B", 6.6, 6.1, ACCENT_TEAL, RGBColor(0,0,0))
add_pill(sl, "SQLite", 8.3, 6.1, RGBColor(0x78, 0x35, 0x00))
add_pill(sl, "Python", 9.3, 6.1, RGBColor(0xF5, 0x9E, 0x0B), RGBColor(0,0,0))


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
slide_header(sl, "Problem Statement", "The Challenge & The Objective")

challenges = [
    ("No Personalisation",     "Generic platforms give random questions — not tailored to the candidate's resume, skills, or target role."),
    ("Zero Real-Time Feedback","Candidates answer into a void. No scoring, no depth evaluation, no actionable improvement tips."),
    ("Static Difficulty",      "A Junior developer and a Senior Engineer receive identical questions — practice is either too easy or overwhelming."),
    ("No Progress Tracking",   "Without historical data, a candidate cannot see whether they are improving session to session."),
]

cols = [(0.3, 1.25), (6.85, 1.25), (0.3, 4.0), (6.85, 4.0)]
for (l, t), (title, desc) in zip(cols, challenges):
    add_rect(sl, l, t, 6.2, 2.55, fill=BG_CARD, line=ACCENT_PURP)
    add_text(sl, "⚠  " + title, l+0.15, t+0.15, 5.9, 0.4, size=15, bold=True, color=ACCENT_RED)
    add_text(sl, desc, l+0.15, t+0.65, 5.9, 1.75, size=13, color=TEXT_WHITE)

# Objective banner
add_rect(sl, 0.3, 6.65, 12.73, 0.65, fill=RGBColor(0x1E, 0x28, 0x4A), line=ACCENT_BLUE)
add_text(sl, "🎯  Objective: Build an intelligent AI-powered mock interviewer that personalises every question, evaluates answers in real time via IBM Watsonx, and delivers data-driven improvement reports.",
         0.5, 6.72, 12.3, 0.5, size=12, color=ACCENT_BLUE)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — PROPOSED SOLUTION
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
slide_header(sl, "Proposed Solution", "Five Core Mechanisms")

solutions = [
    ("📄  Resume Parsing",       "PDF / DOCX file uploaded by candidate is parsed server-side into plain text and injected into every AI prompt for hyper-personalised questions.", ACCENT_TEAL),
    ("🔍  RAG Pipeline",         "ChromaDB vector database stores 500+ curated questions. Top 20 semantically relevant questions are retrieved and used as grounding context for the LLM.", ACCENT_BLUE),
    ("🎚  Adaptive Difficulty",  "SQLite tracks every score. Avg < 5 → easier questions. Avg > 8 → senior-level questions. Difficulty adjusts in real time.", ACCENT_PURP),
    ("🤖  IBM Watsonx Evaluation","Every submitted answer is scored by LLaMA 3.3 70B: Score /10, Strengths, Weaknesses, Improvement Tips, and a rewritten model answer.", RGBColor(0x17, 0x5C, 0xD6)),
    ("📊  Session Reporting",    "On 'End Interview', the AI generates a brief, bulleted final report with average score and a targeted study plan for weak areas.", TEXT_YELLOW),
]

y = 1.25
for title, desc, col in solutions:
    add_rect(sl, 0.3, y, 12.73, 1.0, fill=BG_CARD, line=col)
    add_text(sl, title, 0.5, y+0.05, 3.5, 0.4, size=14, bold=True, color=col)
    add_text(sl, desc, 4.0, y+0.1, 9.0, 0.75, size=13, color=TEXT_WHITE)
    y += 1.12


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
slide_header(sl, "Technology Stack", "Tools, Libraries & Platforms Used")

tech = [
    ("Frontend",       "HTML5 · CSS3 · Vanilla JavaScript",               "Glassmorphic dark UI, animations, REST API calls via fetch()", ACCENT_PURP),
    ("Backend",        "Python · FastAPI · SlowAPI",                       "Async REST API server with rate limiting middleware", ACCENT_BLUE),
    ("AI / LLM",       "IBM Watsonx — meta-llama/llama-3-3-70b-instruct", "Question generation, evaluation, strategy & report via IAM auth", RGBColor(0x17, 0x5C, 0xD6)),
    ("Vector DB",      "ChromaDB + ONNX MiniLM Embeddings",                "Semantic RAG retrieval of 500+ curated interview questions", ACCENT_TEAL),
    ("Session DB",     "SQLite (via Python sqlite3)",                       "Persistent score storage, adaptive difficulty, history page", TEXT_YELLOW),
    ("Resume Parsing", "PyPDF2 · python-docx",                             "Server-side extraction from PDF and DOCX files", ACCENT_TEAL),
    ("Security",       "python-dotenv · CORS Middleware",                  "API keys never hardcoded; cross-origin requests locked down", ACCENT_RED),
    ("Deployment",     "Docker · start_backend.bat",                       "Containerised; single-command local launch", TEXT_MUTED),
]

cols_count = 2
col_w = 6.2
for i, (layer, tools, role, col) in enumerate(tech):
    col_idx = i % 2
    row_idx = i // 2
    l = 0.3 + col_idx * (col_w + 0.5)
    t = 1.25 + row_idx * 1.45
    add_rect(sl, l, t, col_w, 1.3, fill=BG_CARD, line=col)
    add_text(sl, layer, l+0.15, t+0.08, col_w-0.3, 0.35, size=13, bold=True, color=col)
    add_text(sl, tools, l+0.15, t+0.42, col_w-0.3, 0.32, size=11, bold=True, color=TEXT_WHITE)
    add_text(sl, role,  l+0.15, t+0.75, col_w-0.3, 0.48, size=11, color=TEXT_MUTED)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — WORKFLOW DIAGRAM
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
slide_header(sl, "Workflow Diagram", "End-to-End User & System Flow")
sl.shapes.add_picture(WORKFLOW_IMG, Inches(1.5), Inches(1.2), Inches(10.3), Inches(6.0))


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE BLUEPRINT
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
slide_header(sl, "Architecture Blueprint", "3-Tier System Design")
sl.shapes.add_picture(ARCH_IMG, Inches(1.0), Inches(1.2), Inches(11.33), Inches(6.0))


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — ROLE OF SPECIFIC TECH
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
slide_header(sl, "Role of Specific Technologies", "How Each Component Contributes")

roles = [
    ("🧠  IBM Watsonx AI", ACCENT_BLUE,
     ["Brain of the entire application — generates all AI content",
      "Authenticates via IBM IAM token exchange (enterprise-grade security)",
      "Handles: prep strategy, questions, evaluation, final report"]),
    ("🔍  ChromaDB & RAG", ACCENT_TEAL,
     ["Long-term memory — stores 500+ curated expert questions as vector embeddings",
      "Retrieves top 20 semantically relevant questions before prompting the LLM",
      "Grounds LLM output in real interview patterns — reduces hallucinations"]),
    ("⚡  FastAPI", ACCENT_PURP,
     ["Async-first — all Watsonx calls are non-blocking (concurrent users supported)",
      "SlowAPI middleware rate-limits each endpoint to prevent API quota abuse",
      "Clean REST API consumed by the frontend via fetch()"]),
    ("🗄  SQLite", TEXT_YELLOW,
     ["Persists every question-score pair across sessions",
      "Powers the adaptive difficulty engine — query last 3 scores in real time",
      "Enables the Historical Scores page showing all past sessions"]),
]

positions = [(0.3, 1.25), (6.85, 1.25), (0.3, 4.0), (6.85, 4.0)]
for (l, t), (title, col, bullets) in zip(positions, roles):
    add_rect(sl, l, t, 6.2, 3.1, fill=BG_CARD, line=col)
    add_text(sl, title, l+0.15, t+0.1, 5.9, 0.4, size=14, bold=True, color=col)
    by = t + 0.6
    for b in bullets:
        add_text(sl, f"•  {b}", l+0.2, by, 5.8, 0.42, size=12, color=TEXT_WHITE)
        by += 0.45


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — NOVELTY & UNIQUENESS
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
slide_header(sl, "Novelty & Uniqueness", "How We Stand Apart From Common Platforms")

headers = ["Feature", "Common Apps", "AI Interview Trainer"]
col_widths = [3.8, 2.8, 5.4]
col_starts = [0.3, 4.2, 7.1]
row_h = 0.52
data = [
    ["Question Personalisation", "❌ Generic",             "✅ Resume + Role + Experience filtered"],
    ["Adaptive Difficulty",       "❌ Static",              "✅ Adjusts live based on score history"],
    ["Evaluation Depth",          "❌ None / basic",        "✅ Score, Strengths, Tips, Model Answer Rewrite"],
    ["Resume Integration",        "❌ Not supported",       "✅ PDF/DOCX parsed & fed into every AI prompt"],
    ["IBM Cloud Native",          "❌ Third-party LLMs",   "✅ Watsonx IAM auth + enterprise rate limiting"],
    ["Score History",             "❌ Session only",        "✅ SQLite DB + dedicated History page"],
    ["RAG Architecture",          "❌ Pure LLM",            "✅ Hybrid RAG — grounded in curated dataset"],
    ["Final Session Report",      "❌ Not available",       "✅ AI-generated study plan for weak areas"],
]

# Header row
t = 1.25
for ci, (header, cw, cs) in enumerate(zip(headers, col_widths, col_starts)):
    fill = ACCENT_PURP if ci == 2 else RGBColor(0x1E, 0x28, 0x4A)
    add_rect(sl, cs, t, cw, 0.48, fill=fill)
    add_text(sl, header, cs+0.1, t+0.08, cw-0.2, 0.35, size=13, bold=True,
             color=TEXT_WHITE, align=PP_ALIGN.CENTER)

t += 0.52
for ri, row in enumerate(data):
    bg = RGBColor(0x16, 0x1E, 0x35) if ri % 2 == 0 else RGBColor(0x1A, 0x24, 0x40)
    for ci, (cell, cw, cs) in enumerate(zip(row, col_widths, col_starts)):
        add_rect(sl, cs, t, cw, row_h, fill=bg)
        col_c = ACCENT_TEAL if ci == 2 else TEXT_WHITE
        add_text(sl, cell, cs+0.1, t+0.07, cw-0.2, row_h-0.05, size=11.5,
                 color=col_c, align=PP_ALIGN.LEFT)
    t += row_h


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — FUTURE SCOPE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
slide_header(sl, "Future Scope", "What's Next for the AI Interview Trainer")

future = [
    ("🎤", "Voice Interview Mode",      "Web Speech API integration — candidates answer verbally; speech-to-text pipes into the evaluation engine.", ACCENT_PURP),
    ("🌐", "Multi-Language Support",   "Leverage IBM Granite multilingual models to conduct interviews in Hindi, German, French and more.", ACCENT_BLUE),
    ("🧠", "Cross-Session Memory",     "Store past context in ChromaDB so the AI remembers previous weaknesses and focuses on them next session.", ACCENT_TEAL),
    ("📊", "Analytics Dashboard",      "Score trend graphs, topic-wise heatmaps, and improvement velocity charts over weeks.", TEXT_YELLOW),
    ("🏢", "Company-Specific Mode",    "Upload a job description PDF and the system generates questions targeting that company's known interview style.", ACCENT_RED),
    ("☁️", "Cloud Deployment",          "Dockerise and deploy on IBM Code Engine or Google Cloud Run for public, scalable access.", TEXT_MUTED),
]

cols2 = [(0.3, 1.25), (4.55, 1.25), (8.8, 1.25),
         (0.3, 4.0),  (4.55, 4.0),  (8.8, 4.0)]
for (l, t), (icon, title, desc, col) in zip(cols2, future):
    add_rect(sl, l, t, 4.0, 3.0, fill=BG_CARD, line=col)
    add_text(sl, icon, l+0.15, t+0.12, 0.7, 0.55, size=26)
    add_text(sl, title, l+0.15, t+0.7, 3.7, 0.4, size=13, bold=True, color=col)
    add_text(sl, desc,  l+0.15, t+1.15, 3.7, 1.7, size=11.5, color=TEXT_WHITE)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — THANK YOU
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_bg(sl)
add_rect(sl, 0, 2.5, 13.33, 2.5, fill=RGBColor(0x1E, 0x12, 0x40))

add_text(sl, "Thank You", 0, 2.7, 13.33, 1.0,
         size=52, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text(sl, "AI Interview Trainer  ·  IBM Edunet Summer Internship",
         0, 3.8, 13.33, 0.5, size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

add_text(sl, "Built with  🤖  IBM Watsonx  ·  FastAPI  ·  ChromaDB  ·  Python",
         0, 5.8, 13.33, 0.5, size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")